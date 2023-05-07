# training_thread/powerpoint.py

import os
from pptx import Presentation
from pptx.util import Inches


class TrainingPowerpoint:
    def __init__(self, template_path=None):
        if template_path and os.path.exists(template_path):
            self.presentation = Presentation(template_path)
        else:
            self.presentation = Presentation()

    def add_slide(self, title, bullet_points):
        """
        Add a new slide with a title and bullet points to the presentation.

        Args:
            title (str): The title of the new slide.
            bullet_points (list): A list of strings to be added as bullet points to the slide.
        """
        slide_layout = self.presentation.slide_layouts[1]
        slide = self.presentation.slides.add_slide(slide_layout)
        title_placeholder = slide.placeholders[0]
        body_placeholder = slide.placeholders[1]

        title_placeholder.text = title

        bullets = body_placeholder.text_frame
        for point in bullet_points:
            p = bullets.add_paragraph()
            p.text = point
            p.space_before = Inches(0.05)
            p.level = 0

    def save_presentation(self, file_path):
        """
        Save the presentation to a specified file path.

        Args:
            file_path (str): The file path where the presentation should be saved.
        """
        self.presentation.save(file_path)

